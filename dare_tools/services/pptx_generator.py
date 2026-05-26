"""Generate PowerPoint files from DARE presentation artifact JSON."""

from __future__ import annotations

from io import BytesIO
from typing import Any, Dict, Iterable, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

DEFAULT_THEME = {
    "primaryColor": "#2563EB",
    "accentColor": "#14B8A6",
    "backgroundColor": "#F8FAFC",
    "textColor": "#0F172A",
    "mutedTextColor": "#475569",
    "fontFamily": "Aptos",
}


def generate_pptx_bytes(config: Dict[str, Any]) -> bytes:
    """Build a styled PPTX binary from a validated DARE presentation config."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    theme = {**DEFAULT_THEME, **(config.get("theme") or {})}
    blank_layout = prs.slide_layouts[6]

    slides = config.get("slides") or []
    for index, slide_config in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)
        _apply_background(slide, theme["backgroundColor"])
        _add_footer(slide, index + 1, len(slides), theme)

        layout = slide_config.get("layout") or "bullets"
        if layout == "title":
            _render_title_slide(slide, slide_config, theme)
        elif layout == "section":
            _render_section_slide(slide, slide_config, theme)
        elif layout == "twoColumn":
            _render_two_column_slide(slide, slide_config, theme)
        elif layout == "table":
            _render_table_slide(slide, slide_config, theme)
        elif layout == "quote":
            _render_quote_slide(slide, slide_config, theme)
        elif layout == "summary":
            _render_summary_slide(slide, slide_config, theme)
        else:
            _render_bullets_slide(slide, slide_config, theme)

        _add_speaker_notes(slide, slide_config.get("speakerNotes"))

    stream = BytesIO()
    prs.save(stream)
    return stream.getvalue()


def _render_title_slide(slide, slide_config: Dict[str, Any], theme: Dict[str, str]):
    _add_accent_bar(slide, theme)
    _add_textbox(
        slide,
        slide_config.get("title", "Presentation"),
        0.85,
        1.65,
        8.8,
        1.25,
        theme,
        size=38,
        bold=True,
        color=theme["textColor"],
    )
    subtitle = slide_config.get("subtitle")
    if subtitle:
        _add_textbox(
            slide,
            subtitle,
            0.9,
            3.05,
            7.9,
            0.8,
            theme,
            size=18,
            color=theme["mutedTextColor"],
        )
    _add_side_panel(slide, theme)


def _render_section_slide(slide, slide_config: Dict[str, Any], theme: Dict[str, str]):
    _add_accent_bar(slide, theme, height=0.18)
    _add_textbox(
        slide,
        slide_config.get("title", "Section"),
        0.85,
        2.3,
        9.3,
        1.0,
        theme,
        size=34,
        bold=True,
        color=theme["textColor"],
    )
    body = slide_config.get("body") or slide_config.get("subtitle")
    if body:
        _add_textbox(
            slide,
            body,
            0.9,
            3.4,
            8.4,
            1.2,
            theme,
            size=17,
            color=theme["mutedTextColor"],
        )


def _render_bullets_slide(slide, slide_config: Dict[str, Any], theme: Dict[str, str]):
    _add_slide_title(slide, slide_config.get("title", "Key Points"), theme)
    bullets = _clean_strings(slide_config.get("bullets") or [])
    if slide_config.get("body"):
        _add_textbox(
            slide,
            slide_config["body"],
            0.85,
            1.4,
            11.4,
            0.65,
            theme,
            size=15,
            color=theme["mutedTextColor"],
        )
        top = 2.15
    else:
        top = 1.65
    _add_bullets(
        slide,
        bullets,
        1.05,
        top,
        10.9,
        4.7,
        theme,
        size=_fit_bullet_size(bullets, 15),
    )


def _render_two_column_slide(
    slide, slide_config: Dict[str, Any], theme: Dict[str, str]
):
    _add_slide_title(slide, slide_config.get("title", "Comparison"), theme)
    _add_card(
        slide,
        0.85,
        1.55,
        5.55,
        4.95,
        slide_config.get("leftTitle", "Column 1"),
        slide_config.get("leftBullets") or [],
        theme,
    )
    _add_card(
        slide,
        6.85,
        1.55,
        5.55,
        4.95,
        slide_config.get("rightTitle", "Column 2"),
        slide_config.get("rightBullets") or [],
        theme,
    )


def _render_table_slide(slide, slide_config: Dict[str, Any], theme: Dict[str, str]):
    _add_slide_title(slide, slide_config.get("title", "Table"), theme)
    headers = _clean_strings(slide_config.get("headers") or [])
    rows = slide_config.get("rows") or []
    if not headers:
        _add_textbox(slide, "No table data provided.", 1, 2, 10, 0.5, theme)
        return

    row_count = min(len(rows), 8) + 1
    col_count = len(headers)
    table_shape = slide.shapes.add_table(
        row_count, col_count, Inches(0.85), Inches(1.55), Inches(11.65), Inches(4.85)
    )
    table = table_shape.table

    for column_index, header in enumerate(headers):
        cell = table.cell(0, column_index)
        cell.text = header
        _set_cell_fill(cell, theme["primaryColor"])
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.color.rgb = _rgb(theme["backgroundColor"])
                run.font.size = Pt(11)

    for row_index, row in enumerate(rows[:8], start=1):
        values = [str(value) for value in row[:col_count]]
        values += [""] * (col_count - len(values))
        for column_index, value in enumerate(values):
            cell = table.cell(row_index, column_index)
            cell.text = value
            _set_cell_fill(cell, "#FFFFFF" if row_index % 2 else "#EEF2F7")
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
                    run.font.color.rgb = _rgb(theme["textColor"])


def _render_quote_slide(slide, slide_config: Dict[str, Any], theme: Dict[str, str]):
    _add_accent_bar(slide, theme, height=0.18)
    quote = slide_config.get("quote") or slide_config.get("body") or ""
    _add_textbox(
        slide,
        f'"{quote}"',
        1.15,
        1.85,
        10.9,
        2.5,
        theme,
        size=28,
        bold=True,
        color=theme["textColor"],
    )
    attribution = slide_config.get("attribution")
    if attribution:
        _add_textbox(
            slide,
            f"- {attribution}",
            1.2,
            4.65,
            8.5,
            0.5,
            theme,
            size=16,
            color=theme["mutedTextColor"],
        )


def _render_summary_slide(slide, slide_config: Dict[str, Any], theme: Dict[str, str]):
    _add_slide_title(slide, slide_config.get("title", "Summary"), theme)
    bullets = _clean_strings(slide_config.get("bullets") or [])
    box_width = 3.55
    for index, bullet in enumerate(bullets[:3]):
        left = 0.85 + index * 4.05
        _add_numbered_takeaway(slide, index + 1, bullet, left, 2.0, box_width, theme)
    if len(bullets) > 3:
        remaining = bullets[3:]
        _add_bullets(
            slide,
            remaining,
            1.05,
            5.2,
            10.9,
            1.0,
            theme,
            size=_fit_bullet_size(remaining, 13),
        )


def _add_slide_title(slide, title: str, theme: Dict[str, str]):
    _add_accent_bar(slide, theme, height=0.12)
    _add_textbox(
        slide,
        title,
        0.85,
        0.45,
        10.8,
        0.7,
        theme,
        size=26,
        bold=True,
        color=theme["textColor"],
    )


def _add_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    bullets: Iterable[str],
    theme: Dict[str, str],
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb("#FFFFFF")
    shape.line.color.rgb = _rgb("#DDE3EA")
    _add_textbox(
        slide,
        title,
        left + 0.28,
        top + 0.25,
        width - 0.56,
        0.5,
        theme,
        size=16,
        bold=True,
    )
    _add_bullets(
        slide,
        _clean_strings(bullets),
        left + 0.38,
        top + 0.95,
        width - 0.76,
        height - 1.25,
        theme,
        size=_fit_bullet_size(_clean_strings(bullets), 13),
    )


def _add_numbered_takeaway(
    slide,
    number: int,
    text: str,
    left: float,
    top: float,
    width: float,
    theme: Dict[str, str],
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(2.35),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb("#FFFFFF")
    shape.line.color.rgb = _rgb("#DDE3EA")
    _add_textbox(
        slide,
        f"{number:02d}",
        left + 0.3,
        top + 0.25,
        0.8,
        0.4,
        theme,
        size=13,
        bold=True,
        color=theme["primaryColor"],
    )
    _add_textbox(
        slide, text, left + 0.3, top + 0.85, width - 0.6, 1.1, theme, size=15, bold=True
    )


def _add_bullets(
    slide,
    bullets: List[str],
    left: float,
    top: float,
    width: float,
    height: float,
    theme: Dict[str, str],
    size: int = 15,
):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, bullet in enumerate(bullets or [""]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(size)
        paragraph.font.name = theme["fontFamily"]
        paragraph.font.color.rgb = _rgb(theme["textColor"])
        paragraph.space_after = Pt(9)


def _add_textbox(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    theme: Dict[str, str],
    size: int = 14,
    bold: bool = False,
    color: Optional[str] = None,
):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = str(text)
    run.font.name = theme["fontFamily"]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color or theme["textColor"])
    return box


def _add_accent_bar(slide, theme: Dict[str, str], height: float = 0.12):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(theme["primaryColor"])
    shape.line.fill.background()


def _add_side_panel(slide, theme: Dict[str, str]):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(10.85), Inches(0), Inches(2.5), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(theme["primaryColor"])
    shape.line.fill.background()
    accent = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(9.9), Inches(5.6), Inches(1.65), Inches(1.65)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = _rgb(theme["accentColor"])
    accent.line.fill.background()


def _add_footer(slide, current: int, total: int, theme: Dict[str, str]):
    _add_textbox(
        slide,
        f"{current} / {total}",
        11.75,
        7.02,
        0.7,
        0.25,
        theme,
        size=9,
        color=theme["mutedTextColor"],
    )


def _apply_background(slide, color: str):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _add_speaker_notes(slide, notes: Optional[str]):
    if not notes:
        return
    try:
        text_frame = slide.notes_slide.notes_text_frame
        text_frame.text = str(notes)
    except Exception:
        return


def _set_cell_fill(cell, color: str):
    cell.fill.solid()
    cell.fill.fore_color.rgb = _rgb(color)


def _rgb(value: str):
    value = (value or "#000000").strip().lstrip("#")
    if len(value) == 3:
        value = "".join(char * 2 for char in value)
    try:
        red, green, blue = (
            int(value[0:2], 16),
            int(value[2:4], 16),
            int(value[4:6], 16),
        )
    except Exception:
        red, green, blue = (0, 0, 0)
    return RGBColor(red, green, blue)


def _clean_strings(values: Iterable[Any]) -> List[str]:
    return [str(value).strip() for value in values if str(value).strip()]


def _fit_bullet_size(bullets: List[str], base_size: int) -> int:
    longest = max([len(bullet) for bullet in bullets] or [0])
    if len(bullets) > 5 or longest > 110:
        return max(base_size - 3, 10)
    if len(bullets) > 4 or longest > 85:
        return max(base_size - 2, 10)
    if longest > 65:
        return max(base_size - 1, 10)
    return base_size
